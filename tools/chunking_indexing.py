import os
import shutil
import uuid
import hashlib
import pickle

import fitz  # PyMuPDF
import faiss
from sentence_transformers import SentenceTransformer
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from pptx import Presentation
from docx import Document  # For DOCX support

# =====================================================
# CONFIG
# =====================================================
BASE_DIR = "data"
UNCHUNKED_DIR = os.path.join(BASE_DIR, "unchunked")
CHUNKED_DIR = os.path.join(BASE_DIR, "chunked")
FAISS_DIR = os.path.join(BASE_DIR, "faiss")

MAX_TOKENS = 300
OVERLAP = 100
EMBED_DIM = 384  # all-MiniLM-L6-v2

os.makedirs(UNCHUNKED_DIR, exist_ok=True)
os.makedirs(CHUNKED_DIR, exist_ok=True)
os.makedirs(FAISS_DIR, exist_ok=True)

# =====================================================
# LOAD EMBEDDING MODEL
# =====================================================
embedder = SentenceTransformer("all-MiniLM-L6-v2")

# =====================================================
# LOAD OR INIT FAISS
# =====================================================
index_path = os.path.join(FAISS_DIR, "index.faiss")
meta_path = os.path.join(FAISS_DIR, "metadata.pkl")

if os.path.exists(index_path):
    index = faiss.read_index(index_path)
    metadata = pickle.load(open(meta_path, "rb"))
else:
    index = faiss.IndexFlatL2(EMBED_DIM)
    metadata = []

processed_hashes = {m["file_hash"] for m in metadata}

# =====================================================
# UTILS
# =====================================================
def compute_hash(path):
    h = hashlib.md5()
    with open(path, "rb") as f:
        h.update(f.read())
    return h.hexdigest()

# =====================================================
# CONVERT TXT / PPTX / DOCX → PDF
# =====================================================
def convert_to_pdf(path):
    if path.endswith(".pdf"):
        return path

    base = os.path.splitext(os.path.basename(path))[0]
    pdf_path = os.path.join(UNCHUNKED_DIR, base + ".pdf")

    c = canvas.Canvas(pdf_path, pagesize=letter)
    width, height = letter
    y = height - 40

    lines = []

    if path.endswith(".txt"):
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            lines = f.readlines()

    elif path.endswith(".pptx"):
        prs = Presentation(path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    lines.append(shape.text)
    
    elif path.endswith(".docx"):
        doc = Document(path)
        for para in doc.paragraphs:
            if para.text.strip():
                lines.append(para.text)

    for line in lines:
        for subline in line.split("\n"):
            c.drawString(40, y, subline[:1000])
            y -= 14
            if y < 40:
                c.showPage()
                y = height - 40

    c.save()
    return pdf_path

# =====================================================
# EXTRACT TEXT FROM PDF
# =====================================================
def extract_text(pdf_path):
    doc = fitz.open(pdf_path)
    pages = [page.get_text("text") for page in doc]
    return "\n".join(pages)

# =====================================================
# PARAGRAPH-AWARE TOKEN CHUNKING
# =====================================================
def chunk_text(text):
    chunks = []
    paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]

    for para in paragraphs:
        tokens = para.split()

        if len(tokens) <= MAX_TOKENS:
            chunks.append(" ".join(tokens))
            continue

        start = 0
        while start < len(tokens):
            end = start + MAX_TOKENS
            chunks.append(" ".join(tokens[start:end]))
            start += MAX_TOKENS - OVERLAP

    return chunks

# =====================================================
# MAIN PIPELINE
# =====================================================
for filename in os.listdir(UNCHUNKED_DIR):
    if not filename.lower().endswith((".pdf", ".txt", ".pptx", ".docx")):
        continue

    file_path = os.path.join(UNCHUNKED_DIR, filename)
    print(f"\n📄 Processing: {filename}")

    # Convert if needed
    pdf_path = convert_to_pdf(file_path)

    # Compute hash AFTER conversion
    file_hash = compute_hash(pdf_path)

    if file_hash in processed_hashes:
        print("⏭ Already indexed — skipping")
        continue

    # Extract + chunk
    text = extract_text(pdf_path)
    chunks = chunk_text(text)

    if not chunks:
        print("⚠ No text found — skipping")
        continue

    # Embed + store
    embeddings = embedder.encode(chunks, convert_to_numpy=True)
    index.add(embeddings)

    for i, chunk in enumerate(chunks):
        metadata.append({
            "id": str(uuid.uuid4()),
            "file": os.path.basename(pdf_path),
            "chunk_id": i,
            "file_hash": file_hash,
            "text": chunk
        })

    # Move PDF to chunked folder
    shutil.move(pdf_path, os.path.join(CHUNKED_DIR, os.path.basename(pdf_path)))
    processed_hashes.add(file_hash)

    print(f" Chunked {len(chunks)} chunks")

# =====================================================
# SAVE FAISS + METADATA
# =====================================================
faiss.write_index(index, index_path)
pickle.dump(metadata, open(meta_path, "wb"))

print("\n🎉 All files processed successfully")
